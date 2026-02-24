import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

# --- Custom Logger Import ---
from lib.helper.Logger import LoggerSetup 

class PresentationCreatorWithDetails:
    def __init__(self, template_path, output_path, author_name, title, author_pos, title_pos, progress_callback=None):
        # --- Initialize Logger ---
        self.logger = LoggerSetup(
            log_name='Cal PowerPoint', 
            log_dir_relative_path=r"output\Log"
        ).get_logger()
        
        self.logger.info("Initializing PresentationCreatorWithDetails...")

        try:
            self.template_path = template_path
            self.output_path = output_path
            self.author_name = author_name
            self.title = title
            self.author_pos = author_pos
            self.title_pos = title_pos
            self.progress_callback = progress_callback
            
            if not os.path.exists(self.template_path):
                raise FileNotFoundError(f"Template file not found: {self.template_path}")
                
            self.prs = self.load_template()
            self.logger.info("Template loaded successfully.")

        except Exception as e:
            self.logger.critical(f"Initialization failed: {e}")
            raise

    def load_template(self):
        return Presentation(self.template_path)

    def remove_unused_placeholders(self, slide):
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder or (shape.has_text_frame and not shape.text_frame.text.strip()):
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def add_author_and_title(self):
        try:
            slide = self.prs.slides[0]
            self.remove_unused_placeholders(slide)
            
            # Add author name
            author_shape = slide.shapes.add_textbox(Inches(self.author_pos[0]), Inches(self.author_pos[1]), Inches(8), Inches(1))
            author_text_frame = author_shape.text_frame
            author_p = author_text_frame.paragraphs[0]
            author_p.text = self.author_name
            author_p.font.size = Pt(24)
            author_p.font.bold = True
            author_p.font.color.rgb = RGBColor(0, 0, 0)

            # Add title
            title_shape = slide.shapes.add_textbox(Inches(self.title_pos[0]), Inches(self.title_pos[1]), Inches(8), Inches(1))
            title_text_frame = title_shape.text_frame
            title_p = title_text_frame.paragraphs[0]
            title_p.text = self.title
            title_p.font.size = Pt(30)
            title_p.font.bold = True
            title_p.font.color.rgb = RGBColor(0, 0, 0)
        except Exception as e:
            self.logger.error(f"Error adding author/title: {e}")

    def get_positions(self):
        # Updated position mapping for 4 subplots
        return [
            (0.5, 1),   # Top-left
            (7, 1),     # Top-right
            (0.5, 4),   # Bottom-left
            (7, 4)      # Bottom-right
        ]

    def group_by_ch(self, path1):
        # Logging start of grouping (High level)
        self.logger.info(f"Grouping files from: {path1}")
        
        grouped_files = {}
        if not os.path.exists(path1):
            self.logger.error(f"Image directory not found: {path1}")
            return grouped_files

        image_files = [file for file in os.listdir(path1) if file.endswith('.png')]
        
        # Inner loop: No logging here for performance
        for file in image_files:
            # Use regex to find CH<number> or CH[<number>]
            match = re.search(r'CH(\d+)|CH\[(\d+)\]', file)
            ch_number_key = "Unknown" 
            
            if match:
                number_str = match.group(1) or match.group(2)
                ch_number_key = f"CH{number_str}"
            
            if ch_number_key not in grouped_files:
                grouped_files[ch_number_key] = []
            grouped_files[ch_number_key].append(file)
            
        self.logger.info(f"Found {len(grouped_files)} unique channel groups.")
        return grouped_files

    def create_subplot_slides(self, path1, path2, batch_df=None):
        self.logger.info("Starting slide creation process...")
        try:
            # Filter batch file if provided (only process rows with Enb Scrip = 'v')
            if batch_df is not None and 'Enb Scrip' in batch_df.columns:
                batch_df = batch_df[batch_df['Enb Scrip'] == 'v']
                self.logger.info(f"Filtered batch to {len(batch_df)} rows with Enb Scrip == 'v'")
            
            grouped_files = self.group_by_ch(path1)

            def extract_ch_number(ch_key):
                match = re.search(r'CH(\d+)', ch_key)
                if match:
                    return int(match.group(1)) 
                else:
                    return float('inf')

            sorted_ch_keys = sorted(grouped_files.keys(), key=extract_ch_number)
            total_groups = len(sorted_ch_keys)

            positions = self.get_positions()
            img_width = 5.3
            img_height = 3

            # Outer Loop: Iterating Groups
            for idx, ch in enumerate(sorted_ch_keys):
                # Report progress dynamically
                if self.progress_callback:
                    progress = int((idx / total_groups) * 100)  # 0% to 100%
                    self.progress_callback(progress)
                image_files = grouped_files[ch]
                images_added = 0
                slide = None

                # Inner Loop: Iterating Images (HEAVY PROCESSING - NO LOGGING)
                for image_file in image_files:
                    if images_added == 0:
                        slide_layout = self.prs.slide_layouts[5]
                        slide = self.prs.slides.add_slide(slide_layout)
                        self.remove_unused_placeholders(slide)

                        # Add CH title to the slide
                        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
                        title_text_frame = title_shape.text_frame
                        title_p = title_text_frame.paragraphs[0]
                        title_p.text = f"Group: {ch}"
                        title_p.font.size = Pt(24)
                        title_p.font.bold = True
                        title_p.font.color.rgb = RGBColor(204, 0, 0)

                    # Get image and corresponding HTML paths
                    image_path = os.path.join(path1, image_file)
                    html_file = os.path.splitext(image_file)[0] + '.html'
                    html_path = os.path.join(path2, html_file)

                    # Determine position
                    left, top = positions[images_added]

                    # Add image
                    if os.path.exists(image_path):
                        # 1. Add image first
                        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(img_width), height=Inches(img_height))
                        
                        # 2. Add title text box
                        title_left = left
                        title_top = top   
                        title_width = img_width
                        title_height = 0.3 

                        title_shape = slide.shapes.add_textbox(Inches(title_left), Inches(title_top), Inches(title_width), Inches(title_height))

                        # 3. Set textbox background fill to white
                        fill = title_shape.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(255, 255, 255) 

                        # 4. Format the text
                        text_frame = title_shape.text_frame
                        p = text_frame.paragraphs[0]
                        p.text = os.path.splitext(image_file)[0] 
                        p.font.size = Pt(9)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0, 0, 0) 
                        p.alignment = PP_ALIGN.LEFT 

                    # Add HTML link
                    if os.path.exists(html_path):
                        link_left = left + img_width + 0.1
                        link_top = top + (img_height / 2) - 0.25
                        link_textbox = slide.shapes.add_textbox(Inches(link_left), Inches(link_top), Inches(1), Inches(0.5))
                        text_frame = link_textbox.text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = "Html"
                        run.font.size = Pt(12)
                        run.font.color.rgb = RGBColor(0, 0, 255)
                        run.font.underline = True
                        link_textbox.click_action.hyperlink.address = html_path

                    images_added += 1
                    if images_added == 4:
                        images_added = 0

            # Remove second slide if it exists
            slides_element = self.prs.slides._sldIdLst
            if len(slides_element) > 1:
                slide_id = slides_element[1]
                slides_element.remove(slide_id)

            self.add_author_and_title()
            
            # Create output directory if needed
            os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
            
            self.prs.save(self.output_path)
            self.logger.info(f"Presentation saved successfully to: {self.output_path}")
            print(f"Presentation saved to {self.output_path}")

        except Exception as e:
            self.logger.error(f"An error occurred during slide creation: {e}", exc_info=True)
            print(f"An error occurred: {e}")


# Example usage:
if __name__ == "__main__":
    path1 = r"C:\Users\ro898771\Documents\QuickMi2e\output\AutoImage"  # Image directory
    path2 = r"C:\Users\ro898771\Documents\QuickMi2e\output\AutoHtml"   # HTML directory
    template_path = r"C:\Users\ro898771\Documents\QuickMi2e\setting\CFG\pptTemplate\Template.pptx"
    
    author_name = "Author Name"
    title = "Trace File"
    FileName ="Hey" 
    output_path = r"C:\Users\ro898771\Documents\QuickMi2e\output\pptx\{file}.pptx".format(file=FileName)
    author_pos = (0.5, 5.2)
    title_pos = (0.5, 4)

    try:
        ppt_creator = PresentationCreatorWithDetails(template_path, output_path, author_name, title, author_pos, title_pos)
        ppt_creator.create_subplot_slides(path1, path2)
    except Exception as main_e:
        print(f"Execution failed: {main_e}")