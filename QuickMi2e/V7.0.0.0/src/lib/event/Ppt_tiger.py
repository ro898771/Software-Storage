import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

# --- Custom Logger Import ---
from lib.helper.Logger import LoggerSetup 

class PresentationCreator:
    def __init__(self, image_folder, template_path, output_path, author_name, title, author_pos, title_pos, df, progress_callback=None):
        # --- Initialize Logger ---
        self.logger = LoggerSetup(
            log_name='Tiger-Cntrace PowerPoint', 
            log_dir_relative_path=r"output\Log"
        ).get_logger()

        self.logger.info("Initializing PresentationCreator...")

        self.image_folder = image_folder
        self.template_path = template_path
        self.output_path = output_path
        self.author_name = author_name
        self.title = title
        self.author_pos = author_pos
        self.title_pos = title_pos
        self.df = df
        self.progress_callback = progress_callback
        
        try:
            self.prs = self.load_presentation()
            self.logger.info(f"Template loaded successfully from: {self.template_path}")
        except Exception as e:
            self.logger.critical(f"Failed to load PPTX template: {e}")
            raise
    
    def load_presentation(self):
        return Presentation(self.template_path)

    def remove_slide(self, slide):
        try:
            slides = list(self.prs.slides._sldIdLst)
            for sldId in slides:
                if sldId.rId == slide.slide_id:
                    self.prs.slides._sldIdLst.remove(sldId)
                    break
        except Exception as e:
            self.logger.error(f"Error removing slide: {e}")

    def add_author_title_to_first_slide(self):
        try:
            slide = self.prs.slides[0]
            self.remove_unused_placeholders(slide)
            # Add author name
            author_shape = slide.shapes.add_textbox(Inches(self.author_pos[0]), Inches(self.author_pos[1]), Inches(8), Inches(1))
            text_frame = author_shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = self.author_name
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT

            # Add title
            title_shape = slide.shapes.add_textbox(Inches(self.title_pos[0]), Inches(self.title_pos[1]), Inches(8), Inches(1))
            text_frame = title_shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = self.title
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            
        except Exception as e:
            self.logger.error(f"Error adding author/title to first slide: {e}")

    def create_image_slides(self):
        self.logger.info("Starting slide creation process...")
        
        # Filter to process only enabled scripts ("v")
        if 'Enb Scrip' in self.df.columns:
            self.df = self.df[self.df['Enb Scrip'] == 'v']
            self.logger.info(f"Filtered to {len(self.df)} rows with Enb Scrip == 'v'")
        
        # Add author and title to the first slide
        self.add_author_title_to_first_slide()

        # Sort group names by CH\d+ pattern
        def extract_ch_number(name):
            match = re.search(r'CH(\d+)', str(name)) # Use str() for safety
            return int(match.group(1)) if match else float('inf')  # Non-matching names go last

        sorted_group_names = sorted(self.df['Group Name'].unique(), key=extract_ch_number)
        
        self.logger.info(f"Found {len(sorted_group_names)} groups to process.")

        slides_to_remove = []
        total_groups = len(sorted_group_names)

        # [PERFORMANCE NOTE] Logging inside this loop is minimized to avoid slowdowns
        for idx, group_name in enumerate(sorted_group_names):
            # Report progress dynamically
            if self.progress_callback:
                progress = int((idx / total_groups) * 100)  # 0% to 100%
                self.progress_callback(progress)
            # Use .copy() to avoid SettingWithCopyWarning
            group_df = self.df[self.df['Group Name'] == group_name].copy() 

            # --- START: MODIFIED SORTING LOGIC ---
            
            # Helper function to get CH number from title
            def extract_ch_from_title(title):
                match = re.search(r'CH(\d+)', str(title))
                return int(match.group(1)) if match else float('inf')

            # Helper function to get # number from title
            def extract_hash_from_title(title):
                match = re.search(r'#(\d+)', str(title))
                return int(match.group(1)) if match else float('inf')

            # 1. Create temporary sorting columns
            group_df['_sort_ch'] = group_df['Plot Title'].apply(extract_ch_from_title)
            group_df['_sort_hash'] = group_df['Plot Title'].apply(extract_hash_from_title)

            # 2. Sort by CH number, then by # number
            group_df = group_df.sort_values(by=['_sort_ch', '_sort_hash'])

            # 3. Drop temporary columns (optional, but good practice)
            group_df = group_df.drop(columns=['_sort_ch', '_sort_hash'])
            
            # --- END: MODIFIED SORTING LOGIC ---

            # Fetch the number of plots and positions
            # Ensure group_df is not empty after filtering
            if group_df.empty:
                continue
                
            num_plots = group_df['# Num Plot'].iloc[0]
            pos_locations = group_df['# Pos Location'].tolist()
            positions = self.get_positions(pos_locations)

            # Create a new slide for each group
            try:
                slide_layout = self.prs.slide_layouts[5]  # Adjust layout index if needed
                slide = self.prs.slides.add_slide(slide_layout)
                self.remove_unused_placeholders(slide)

                # Set the slide title to the 'Group Name'
                if slide.shapes.title:
                    slide.shapes.title.text = group_name
                else:
                    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
                    text_frame = title_shape.text_frame
                    p = text_frame.paragraphs[0]
                    p.text = group_name
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.LEFT
            except Exception as e:
                self.logger.error(f"Error creating slide for group {group_name}: {e}")
                continue

            # Get sorted 'Plot Title's and map to image filenames
            plot_titles = group_df['Plot Title'].tolist()
            images = [os.path.join(self.image_folder, f"{self.sanitize_filename(title)}.png") for title in plot_titles]

            # Determine image dimensions
            img_width, img_height = self.get_image_dimensions(num_plots)

            # Add images and HTML links to the slide
            images_added = self.add_images_to_slide(slide, images, img_width, img_height, positions, plot_titles, num_plots)

            # Track slide for removal if no images were added
            if not images_added:
                slides_to_remove.append(slide)

        # --- FIX: Remove empty slides correctly ---
        # Iterate over the list of slides to remove and delete them
        for slide in slides_to_remove:
            self.remove_slide(slide)

        # Optional: remove second slide if needed
        try:
            slides_element = self.prs.slides._sldIdLst
            if len(slides_element) > 1:
                slide_id = slides_element[1]
                slides_element.remove(slide_id)
        except Exception as e:
            self.logger.warning(f"Could not remove second slide (non-critical): {e}")

        # Save the presentation
        try:
            self.prs.save(self.output_path)
            self.logger.info(f"Presentation successfully saved to: {self.output_path}")
        except Exception as e:
            self.logger.critical(f"Failed to save presentation: {e}")

    def get_positions(self, pos_locations):
        # Mapping positions based on '# Pos Location'
        position_mapping = {
            1: (0.5, 1),   # Top-left
            2: (7, 1),     # Top-right
            3: (0.5, 4),   # Bottom-left
            4: (7, 4)      # Bottom-right
        }
        # Handle cases where pos might not be in the map, default to top-left
        positions = [position_mapping.get(pos, (0.5, 1)) for pos in pos_locations]
        return positions

    def get_image_dimensions(self, num_plots):
        if num_plots > 1:
            img_width = 5.3
            img_height = 3
        else:
            # Make single plot larger
            img_width = 10
            img_height = 5.9
        return img_width, img_height


    def add_images_to_slide(self, slide, images, img_width, img_height, positions, plot_titles, num_plots):
            images_added = False
            for i, image_path in enumerate(images):
                if i < len(positions) and i < len(plot_titles): # Ensure we don't go out of bounds
                    left, top = positions[i]
                    if os.path.exists(image_path):
                        try:
                            images_added = True

                            # --- 1. ADD IMAGE FIRST (so it's in the back) ---
                            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(img_width), height=Inches(img_height))

                            # --- 2. ADD TITLE TEXTBOX (on top of the image) ---
                            title_left = left
                            # Adjust title position: move down by 0.6 inches (6 steps) when only 1 plot
                            title_top = top + 0.2 if num_plots == 1 else top
                            title_width = img_width
                            title_height = 0.3 # Height of the title bar

                            title_shape = slide.shapes.add_textbox(Inches(title_left), Inches(title_top), Inches(title_width), Inches(title_height))

                            # --- 3. SET TEXTBOX BACKGROUND FILL TO WHITE ---
                            fill = title_shape.fill
                            fill.solid()
                            fill.fore_color.rgb = RGBColor(255, 255, 255) # White fill

                            # --- 4. FORMAT THE TEXT (Black, Bold, Centered) ---
                            text_frame = title_shape.text_frame
                            p = text_frame.paragraphs[0]
                            p.text = plot_titles[i]
                            # Use Pt(12) for single plot, Pt(9) for multiple plots
                            p.font.size = Pt(12) if num_plots == 1 else Pt(9)
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(0, 0, 0) # Black text
                            p.alignment = PP_ALIGN.CENTER # Center the title

                            # --- 5. ADD HYPERLINK (position remains the same) ---
                            link_left = left + img_width + 0.1
                            link_top = top + (img_height / 2) - 0.25
                            link_textbox = slide.shapes.add_textbox(Inches(link_left), Inches(link_top), Inches(1), Inches(0.5))
                            text_frame = link_textbox.text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = "Html"
                            run.font.size = Pt(12)
                            run.font.color.rgb = RGBColor(0, 0, 255)
                            run.font.underline = True

                            # Create the hyperlink
                            html_file = f"{self.sanitize_filename(plot_titles[i])}.html"
                            # --- FIX: Hardcoded path changed to be relative ---
                            html_folder = os.path.join(os.getcwd(), "output", "AutoHtml")
                            link_path = os.path.join(html_folder, html_file)
                            
                            # Check if HTML file exists before linking
                            if os.path.exists(link_path):
                                link_textbox.click_action.hyperlink.address = link_path
                            else:
                                # Optional: remove the "Html" text if link doesn't exist
                                run.text = "" 
                                # Using debug instead of warning to reduce log noise if many HTMLs are missing
                                self.logger.debug(f"HTML file not found - {link_path}")
                        
                        except Exception as e:
                            self.logger.error(f"Error adding image/elements for {image_path}: {e}")

                    else:
                        self.logger.warning(f"Image file not found - {image_path}")
                else:
                    self.logger.warning(f"Mismatch in images/positions/titles for group.")

            return images_added

    def remove_unused_placeholders(self, slide):
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                # Remove placeholders that are not titles (type 1)
                if shape.placeholder_format.type != 1: 
                    shapes_to_remove.append(shape)
            elif shape.has_text_frame:
                # Remove text boxes with empty text
                text_content = shape.text_frame.text.strip()
                if not text_content:
                    shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception as e:
                # Non-critical error, debug level only
                self.logger.debug(f"Could not remove shape: {e}")


    def sanitize_filename(self, filename):
        # Remove or replace invalid characters from filename
        invalid_chars = ['*', ':', '/', '\\', '?', '"', '<', '>', '|']
        for char in invalid_chars:
            filename = filename.replace(char, '')
        filename = filename.strip()
        return filename

if __name__ == "__main__":
    # Define your DataFrame (Example, as original CSV path is used)
    # ... (DataFrame data omitted for brevity, logic remains the same) ...
    
    # Use dynamic path based on script location
    base_path = os.getcwd()
    
    csv_path = os.path.join(base_path, "setting", "CFG", "Batch", "Module_20250228_135039.csv")
    if not os.path.exists(csv_path):
        print(f"Error: CSV file not found at {csv_path}")
        # Create an empty DataFrame as a fallback
        df = pd.DataFrame(columns=['Group Name', 'Enb Scrip', '# Num Plot', '# Pos Location', 'Plot Title', 'Filtered Folder'])
    else:
        df = pd.read_csv(csv_path)


    # Set your paths and parameters
    image_folder = os.path.join(base_path, "output", "image")
    template_path = os.path.join(base_path, "setting", "CFG", "pptTemplate", "Template.pptx")
    output_folder = os.path.join(base_path, "setting", "pptx")
    output_path = os.path.join(output_folder, "output.pptx")

    # Ensure output directory exists
    os.makedirs(output_folder, exist_ok=True)
    
    author_name = "Author Name"
    title = "Trace File"
    author_pos = (0.5, 5.2)  # Position of the author name (left, top) in inches
    title_pos = (0.5, 4)   # Position of the title (left, top) in inches

    try:
        # Instantiate the PresentationCreator with your DataFrame
        presentation_creator = PresentationCreator(
            image_folder=image_folder,
            template_path=template_path,
            output_path=output_path,
            author_name=author_name,
            title=title,
            author_pos=author_pos,
            title_pos=title_pos,
            df=df  # Pass the DataFrame to the class
        )

        # Create the image slides
        print("Creating presentation with create_image_slides()...")
        presentation_creator.create_image_slides() 
        print(f"Presentation saved to {output_path}")

    except Exception as e:
        print(f"Main execution failed: {e}")